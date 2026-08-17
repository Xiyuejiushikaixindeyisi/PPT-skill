# -*- coding: utf-8 -*-
"""
dept_pages.py — 部门版式脚手架(封面/目录/节标题/内容页/末尾页)+ 页面家具。

坐标来自部门真实文件线框提取与母版底稿截图量测(references/layout-library.md)。
所有颜色/字号/位置从 tokens.json 取,换主题不用改本文件。

用法:
    from pptx_kit import Deck
    import dept_pages as P
    d = Deck("assets/tokens.json")
    P.cover(d, "主标题", subtitle="汇报场合", footer="XX部门 · 2026年7月")
    P.toc(d, ["议题一", "议题二"])
    P.section(d, 2, ["议题一", "议题二", "议题三"])   # 当前第2节高亮
    s = P.content(d, "结论式页标题", label="分组导语")  # 自动画灰线/角标/页码/logo
    P.end(d)
    d.save("deck.pptx")
"""
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def _pct(d, xp, yp, wp, hp):
    return (d.SW*xp/100, d.SH*yp/100, d.SW*wp/100, d.SH*hp/100)

def _sz(d, key, default):
    try: return d.sz(key)
    except KeyError: return default

def _pageno(d):
    return len(d.prs.slides._sldIdLst)

def frame(d, s, page_no=None):
    """内容页页面家具(母版实测):标题下通栏灰线 + 左上角图标 + 左下页码 + 右下logo。"""
    fr = d.th.layout.get("page_frame", {})
    ln = fr.get("line", {})
    x0 = d.SW*ln.get("x0_pct",5.5)/100; x1 = d.SW*ln.get("x1_pct",94.4)/100
    y  = d.SH*ln.get("y_pct",9.8)/100
    bar = s.shapes.add_shape(1, Inches(x0), Inches(y), Inches(x1-x0), Pt(ln.get("w_pt",0.75)))
    bar.fill.solid(); bar.fill.fore_color.rgb = d.c(ln.get("color","border_gray"))
    bar.line.fill.background(); bar.shadow.inherit = False
    ic = fr.get("icon", {})
    ipath = os.path.join(d.th.dir, ic.get("file","corner-icon.png"))
    if os.path.exists(ipath):
        s.shapes.add_picture(ipath, Inches(d.SW*ic.get("x_pct",3.4)/100),
                             Inches(d.SH*ic.get("y_pct",4.0)/100), height=Inches(ic.get("h_in",0.33)))
    pn = fr.get("page_num", {})
    if page_no is None: page_no = _pageno(d)
    d.text(s, d.SW*pn.get("x_pct",4.4)/100, d.SH*pn.get("y_pct",95.0)/100, 1.2, 0.2,
           [[(f"{pn.get('prefix','page ')}{page_no}", pn.get("sz_pt",8), d.c("muted"), False)]])
    d.logo(s)
    return s

def cover(d, title, subtitle=None, footer=None, accent_bar=True):
    s = d.add_slide(); s.name = "frame:cover"
    if accent_bar:
        d.box(s, d.SW*5.4/100, d.SH*8.8/100 - 0.14, 1.2, 0.06, fill="accent", line=None)
    x,y,w,h = _pct(d, 5.4, 8.8, 80, 39)
    d.text(s, x, y, w, h, [[(title, _sz(d,"cover_title",40), d.c("title"), True)]],
           align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM, ls=1.1)
    if subtitle:
        x,y,w,h = _pct(d, 5.4, 56.4, 60, 5)
        d.text(s, x, y, w, h, [[(subtitle, 14, d.c("text"), False)]])
    if footer:
        x,y,w,h = _pct(d, 5.4, 91.0, 40, 3.5)
        d.text(s, x, y, w, h, [[(footer, _sz(d,"caption",9), d.c("muted"), False)]])
    d.logo(s)
    return s

def toc(d, items, heading="目录"):
    s = content(d, heading); s.name = "frame:toc"
    n = len(items)
    two_col = n > 8
    per_col = (n + 1)//2 if two_col else n
    col_w = (d.SW*89/100 - 0.3)/2 if two_col else d.SW*70/100
    x0 = d.SW*5.5/100; y0 = d.SH*16/100
    row_h = min(0.66, (d.SH*0.78 - y0) / max(per_col,1))
    for i, it in enumerate(items):
        col = i // per_col; row = i % per_col
        x = x0 + col*(col_w + 0.3); y = y0 + row*row_h
        d.box(s, x, y+0.06, 0.45, 0.45, fill="accent", line=None,
              paras=[[(f"{i+1:02d}", 14, d.c("white"), True)]])
        d.text(s, x+0.62, y+0.06, col_w-0.62, 0.45,
               [[(it, 14, d.c("text"), False)]], anchor=MSO_ANCHOR.MIDDLE)
        d.box(s, x+0.62, y+row_h-0.03, col_w-0.62, 0.012, fill="gray_line", line=None)
    return s

def section(d, num, sections, image=None, subnotes=None):
    """章节页(4 文件 ×15 页实测):左 37.5% 藏青栏(可叠图)+ 60pt 白色大序号;
    右侧竖分隔线 + 议程列表,当前节藏青加粗、其余灰色,前置小圆点。
    num: 当前节序号(1 起);sections: 全部节标题列表;image: 左栏图片路径(可选);
    subnotes: 当前节下的 2–3 条子卡片文字(可选)。"""
    s = d.add_slide(); s.name = "frame:section"
    sp = d.th.layout.get("section_panel", {})
    pw = d.SW*sp.get("w_pct",37.5)/100
    d.box(s, 0, 0, pw, d.SH, fill=sp.get("fill","navy"), line=None)
    if image and os.path.exists(image):
        s.shapes.add_picture(image, 0, 0, Inches(pw), Inches(d.SH))
    np_ = sp.get("num_pos", {})
    d.text(s, d.SW*np_.get("x_pct",11.3)/100, d.SH*np_.get("y_pct",43.4)/100,
           d.SW*np_.get("w_pct",17.8)/100, d.SH*np_.get("h_pct",14.8)/100,
           [[(f"{num:02d}", sp.get("num_pt",60), d.c("white"), True)]],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    dx = d.SW*sp.get("divider_x_pct",47.5)/100
    dline = s.shapes.add_shape(1, Inches(dx), Inches(d.SH*0.23), Pt(0.75), Inches(d.SH*0.44))
    dline.fill.solid(); dline.fill.fore_color.rgb = d.c("border_gray")
    dline.line.fill.background(); dline.shadow.inherit = False
    n = len(sections); y0 = d.SH*0.235; step = min(0.9, d.SH*0.44/max(n,1))
    for i, t in enumerate(sections, 1):
        cur = (i == num)
        col = d.c("navy") if cur else d.c("muted")
        y = y0 + (i-1)*step
        dot = s.shapes.add_shape(9, Inches(dx-0.07), Inches(y+0.09), Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = d.c("accent") if cur else d.c("navy")
        dot.line.fill.background(); dot.shadow.inherit = False
        d.text(s, dx+0.25, y, 0.9, 0.34, [[(f"{i:02d}", 20, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
        d.text(s, dx+1.15, y, d.SW-dx-1.5, 0.34, [[(t, 20, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
        if cur and subnotes:
            m = len(subnotes); cw = (d.SW-dx-0.6-0.15*(m-1))/m
            for j, sn in enumerate(subnotes):
                d.box(s, dx+0.25+j*(cw+0.15), y+0.42, cw, 0.5,
                      paras=[[(sn, 12, d.c("muted"), True)]], fill="neutral", line=None)
    d.logo(s)
    return s

def content(d, title, label=None):
    """内容页:标题带 (5.5,1.7) 89×8(20pt 加粗黑,结论式长标题)+ 页面家具
    (灰线/角标/页码/logo)。可选 label = 标题下分组导语标签(9/9 文件惯例)。
    主图区 y≈12%→94%(高度≥70%)由调用方按 motif-library 绘制。"""
    s = d.add_slide()
    x,y,w,h = _pct(d, 5.5, 1.7, 89, 8)
    d.text(s, x, y, w, h, [[(title, _sz(d,"title",20), d.c("title"), True)]],
           align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
    frame(d, s)
    if label:
        d.box(s, d.SW*5.5/100, d.SH*11.2/100, min(0.32+0.17*len(label), 3.6), 0.32,
              paras=[[(label, 12, d.c("text"), True)]], fill="container", line=None,
              align=PP_ALIGN.CENTER)
    return s

def end(d, text="谢  谢", subline=None):
    s = d.add_slide(); s.name = "frame:end"
    x,y,w,h = _pct(d, 31, 34, 38, 24)
    d.text(s, x, y, w, h, [[(text, 40, d.c("title"), True)]],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if subline:
        x,y,w,h = _pct(d, 31, 60, 38, 6)
        d.text(s, x, y, w, h, [[(subline, 14, d.c("muted"), False)]], align=PP_ALIGN.CENTER)
    d.logo(s)
    return s
