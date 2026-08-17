# -*- coding: utf-8 -*-
"""
pptx_kit.py — theme-driven python-pptx toolkit for high-density engineering slides.

Every primitive pulls colors / fonts / sizes from a Theme loaded from tokens.json,
so re-branding = swap the JSON. Nothing here is model-specific; any environment with
python-pptx can run it.

Primitives: title, box, circle, pill, harrow, varrow, chevron, diamond, dashline,
vlabel, mini_table, layer_band, gantt (helper), swimlane_lane, logo.

Usage:
    from pptx_kit import Deck
    d = Deck("assets/tokens.json")
    s = d.add_slide()
    d.title(s, "标题")
    d.box(s, x, y, w, h, paras=[[("文字", d.sz("body"), d.c("text"), True)]])  # all sizes >= sz("min") = 11pt
    ...
    d.save("out.pptx")
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

def _rgb(hexstr): return RGBColor(int(hexstr[0:2],16), int(hexstr[2:4],16), int(hexstr[4:6],16))

class Theme:
    def __init__(self, tokens_path):
        with open(tokens_path, "r", encoding="utf-8") as f:
            self.t = json.load(f)
        self.dir = os.path.dirname(os.path.abspath(tokens_path))
    def c(self, key):
        p = self.t["palette"]
        if key in p: return _rgb(p[key])
        g = self.t.get("gantt_status", {})
        if key in g: return _rgb(g[key])
        raise KeyError(f"color '{key}' not in palette/gantt_status")
    def sz(self, key): return self.t["type_scale_pt"][key]
    @property
    def font(self): return self.t["font"]["family"]
    @property
    def layout(self): return self.t["layout"]
    @property
    def density(self): return self.t["density"]

class Deck:
    def __init__(self, tokens_path):
        self.th = Theme(tokens_path)
        self.prs = Presentation()
        L = self.th.layout
        self.SW = L["slide_w_in"]; self.SH = L["slide_h_in"]
        self.prs.slide_width  = Inches(self.SW)
        self.prs.slide_height = Inches(self.SH)
        self.M = L["margin_in"]
        self.blank = self.prs.slide_layouts[6]

    # ---- passthrough helpers ----
    def c(self, k): return self.th.c(k)
    def sz(self, k): return self.th.sz(k)

    def add_slide(self): return self.prs.slides.add_slide(self.blank)
    def save(self, path): self.prs.save(path)

    # ---- low-level font ----
    def _font(self, run, size, color, bold=False, italic=False):
        name = self.th.font
        run.font.size = Pt(size); run.font.color.rgb = color
        run.font.bold = bold; run.font.italic = italic; run.font.name = name
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {}); rPr.append(el)
            el.set("typeface", name)

    def _noauto(self, tf):
        tf.word_wrap = True
        try: tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception: pass

    def _write(self, tf, paras, align, anchor, ls, sa):
        self._noauto(tf); tf.vertical_anchor = anchor
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.line_spacing = ls
            p.space_after = Pt(sa); p.space_before = Pt(0)
            for seg in para:
                text, size, color, bold, italic = (tuple(seg) + (False, False))[:5]
                r = p.add_run(); r.text = text
                self._font(r, size, color, bold, italic)

    def _shadow(self, sp):
        spPr = sp._element.spPr
        e = spPr.find(qn('a:effectLst'))
        if e is not None: spPr.remove(e)
        eff = spPr.makeelement(qn('a:effectLst'), {})
        o = eff.makeelement(qn('a:outerShdw'),
            {'blurRad':'42000','dist':'19000','dir':'5400000','rotWithShape':'0'})
        col = o.makeelement(qn('a:srgbClr'), {'val': self.th.t["palette"]["title"]})
        a = col.makeelement(qn('a:alpha'), {'val':'12000'})
        col.append(a); o.append(col); eff.append(o); spPr.append(eff)

    # ---- primitives ----
    def text(self, s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.02, sa=1):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
        self._write(tf, paras, align, anchor, ls, sa); return tb

    def box(self, s, x, y, w, h, paras=None, fill="white", line="border_gray", lw=1.0, rad=None,
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, shadow=False,
            ml=0.06, mr=0.06, mt=0.03, mb=0.03, ls=1.02, sa=1):
        """All text boxes are SQUARE-CORNERED rectangles (hard style rule).
        `rad` is accepted but IGNORED — kept only so older call sites don't break."""
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = self.c(fill) if isinstance(fill,str) else fill
        if line is None: sp.line.fill.background()
        else:
            sp.line.color.rgb = self.c(line) if isinstance(line,str) else line
            sp.line.width = Pt(lw)
        sp.shadow.inherit = False
        if shadow: self._shadow(sp)
        if paras:
            tf = sp.text_frame
            tf.margin_left=Inches(ml); tf.margin_right=Inches(mr)
            tf.margin_top=Inches(mt); tf.margin_bottom=Inches(mb)
            self._write(tf, paras, align, anchor, ls, sa)
        return sp

    def circle(self, s, cx, cy, d, fill, glyph, gcolor, gsize, gbold=True):
        sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(fill) if isinstance(fill,str) else fill
        sp.line.fill.background(); sp.shadow.inherit = False
        tf = sp.text_frame; self._noauto(tf)
        for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = glyph
        self._font(r, gsize, self.c(gcolor) if isinstance(gcolor,str) else gcolor, gbold)
        return sp

    def pill(self, s, x, y, w, h, txt, fill, tc, size=11, bold=True):
        # chips are text boxes too -> square-cornered rectangle, no rounding
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(fill) if isinstance(fill,str) else fill
        sp.line.fill.background(); sp.shadow.inherit = False
        tf = sp.text_frame; self._noauto(tf)
        for m in ('margin_top','margin_bottom'): setattr(tf, m, 0)
        tf.margin_left=Inches(0.02); tf.margin_right=Inches(0.02); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        self._font(r, size, self.c(tc) if isinstance(tc,str) else tc, bold)
        return sp

    def harrow(self, s, x, y, w, h, color="accent", left=False):
        shp = MSO_SHAPE.LEFT_ARROW if left else MSO_SHAPE.RIGHT_ARROW
        sp = s.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(color) if isinstance(color,str) else color
        sp.line.fill.background(); sp.shadow.inherit = False
        try: sp.adjustments[0]=0.55; sp.adjustments[1]=0.55
        except Exception: pass
        return sp

    def varrow(self, s, x, y, w, h, color="accent", up=False):
        shp = MSO_SHAPE.UP_ARROW if up else MSO_SHAPE.DOWN_ARROW
        sp = s.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(color) if isinstance(color,str) else color
        sp.line.fill.background(); sp.shadow.inherit = False
        try: sp.adjustments[0]=0.5; sp.adjustments[1]=0.55
        except Exception: pass
        return sp

    def chevron(self, s, x, y, w, h, txt, fill, tc, size=11.5, bold=True):
        sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(fill) if isinstance(fill,str) else fill
        sp.line.fill.background(); sp.shadow.inherit = False
        try: sp.adjustments[0]=0.5
        except Exception: pass
        tf = sp.text_frame; self._noauto(tf)
        for m in ('margin_top','margin_bottom'): setattr(tf, m, 0)
        tf.margin_left=Inches(0.10); tf.margin_right=Inches(0.02); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        self._font(r, size, self.c(tc) if isinstance(tc,str) else tc, bold)
        return sp

    def diamond(self, s, cx, cy, d, fill="milestone"):
        sp = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(cx), Inches(cy), Inches(d), Inches(d))
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(fill) if isinstance(fill,str) else fill
        sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def dashline(self, s, x, y1, y2, color="accent"):
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2))
        ln.line.color.rgb = self.c(color) if isinstance(color,str) else color
        ln.line.width = Pt(1.2); ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        ln.shadow.inherit = False; return ln

    def vlabel(self, s, x, y, w, h, txt, fill, tc, size=11.5, bold=True, line=None):
        paras = [[(ch, size, self.c(tc) if isinstance(tc,str) else tc, bold)] for ch in txt]
        return self.box(s, x, y, w, h, paras, fill=fill, line=line, ls=0.98, sa=0, ml=0, mr=0)

    def mini_table(self, s, x, y, w, heads, rows, col_ws, hh=0.32, rh=0.34, fs=11):
        cx = x
        for hd, cw in zip(heads, col_ws):
            self.box(s, cx, y, cw, hh, [[(hd, fs, self.c("white"), True)]],
                     fill="title", line=None, align=PP_ALIGN.LEFT, ml=0.07)
            cx += cw
        for r, row in enumerate(rows):
            cx = x; fill = "neutral" if r % 2 == 0 else "nested"
            for c_i, (val, cw) in enumerate(zip(row, col_ws)):
                self.box(s, cx, y+hh+r*rh, cw, rh, [[(val, fs, self.c("text"), c_i==0)]],
                         fill=fill, line=None, align=PP_ALIGN.LEFT, ml=0.07, ls=0.98)
                cx += cw
        return hh + rh*len(rows)

    def layer_band(self, s, main_x, main_w, lane_x, lane_w, y, h, fill, label):
        """母题 F 一层:左侧竖排层名 + 全宽色带。"""
        self.vlabel(s, lane_x, y, lane_w, h, label, "title", "white", 11.5)
        self.box(s, main_x, y, main_w, h, None, fill=fill, line=None, shadow=True)

    def logo(self, s):
        Lg = self.th.layout["logo"]
        w = Lg["width_in"]; h = w * Lg["aspect_h"] / Lg["aspect_w"]
        path = os.path.join(self.th.dir, Lg["file"])
        x = self.SW - Lg["pad_in"] - w
        y = self.SH - Lg["pad_in"] - h
        if os.path.exists(path):
            s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

    def chip(self, s, x, y, w, h, text_, size=None, fill="container", tc="text", bold=False):
        """Small rounded tag chip (dept exception: only shapes <=1.1x0.5in may be rounded)."""
        assert w <= 1.1 and h <= 0.5, "chip too big — use d.box (square) for containers"
        sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        try: sp.adjustments[0] = 0.25
        except Exception: pass
        sp.fill.solid(); sp.fill.fore_color.rgb = self.c(fill) if isinstance(fill,str) else fill
        sp.line.fill.background(); sp.shadow.inherit = False
        tf = sp.text_frame
        for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
        self._write(tf, [[(text_, size or self.sz("caption"), self.c(tc) if isinstance(tc,str) else tc, bold)]],
                    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0, 0)
        return sp

    def title(self, s, t):
        tb = self.th.layout.get("title_band")
        if tb:
            x = self.SW*tb["x_pct"]/100; y = self.SH*tb["y_pct"]/100
            w = self.SW*tb["w_pct"]/100; h = self.SH*tb["h_pct"]/100
        else:
            x, y, w, h = self.M, 0.26, self.SW - 2*self.M, 0.55
        self.text(s, x, y, w, h,
                  [[(t, self.sz("title"), self.c("title"), True)]], anchor=MSO_ANCHOR.MIDDLE)
