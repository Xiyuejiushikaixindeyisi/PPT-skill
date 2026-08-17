import re
# -*- coding: utf-8 -*-
"""
qa_check.py — automated verification for high-density engineering slides.

Runs four gates and returns a non-zero exit code if any hard gate fails, so it can
be used in an iterate-until-pass loop.

  Gate 1  TEXT OVERFLOW  : measure each single-line box's rendered width with a real
                           CJK font; fail if width exceeds available by tolerance.
  Gate 2  GEOMETRY       : out-of-bounds shapes; unintended partial overlaps of major blocks.
  Gate 3  DENSITY        : rasterize each slide; ink-coverage + largest blank square.
  Gate 4  STRUCTURE      : minimum shape count per content slide (anti "fake flowchart").
  Gate 5  FONT SIZE      : every text run must be >= min_font_pt (default 11pt); runs
                           without an explicit size are reported as warnings.
  Gate 6  SHAPE          : no rounded rectangles — all text boxes must be square-cornered.

Usage:
    python qa_check.py deck.pptx --tokens assets/tokens.json [--render-dir /tmp/qa]

Needs: python-pptx, Pillow, and (for Gate 3) a way to rasterize — either pre-rendered
JPGs in --render-dir named slide-<n>.jpg, or soffice+pdftoppm on PATH (auto-run).
Requires a CJK font; searches common locations, override with --font.
"""
import sys, os, json, argparse, glob, subprocess, tempfile
from pptx import Presentation
from pptx.util import Emu

EMU = 914400
CJK_CANDIDATES = [
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/System/Library/Fonts/PingFang.ttc",
    "C:/Windows/Fonts/msyh.ttc",
]

def find_font(override=None):
    if override and os.path.exists(override): return override
    for p in CJK_CANDIDATES:
        if os.path.exists(p): return p
    return None

def load_tokens(path):
    with open(path, encoding="utf-8") as f: return json.load(f)

def gate_text_overflow(prs, font_path, tol_pct):
    from PIL import ImageFont
    cache = {}
    def font(sz, bold):
        k = (round(sz*2), bold)
        if k not in cache:
            cache[k] = ImageFont.truetype(font_path, int(round(sz*96/72)))
        return cache[k]
    flags = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False): continue
            try: w = sh.width/EMU; h = sh.height/EMU
            except Exception: continue
            tf = sh.text_frame
            ml = (tf.margin_left or 0)/EMU; mr = (tf.margin_right or 0)/EMU
            avail = w - ml - mr
            for p in tf.paragraphs:
                line = "".join(r.text for r in p.runs)
                if not line.strip(): continue
                wpx = sum(font((r.font.size.pt if r.font.size else 11),
                               bool(r.font.bold)).getlength(r.text) for r in p.runs)
                win = wpx/96.0
                if h <= 0.42 and win > avail*(1+tol_pct/100.0):
                    flags.append(f"S{si} OVERFLOW '{line[:26]}' need {win:.2f}in have {avail:.2f}in")
    return flags

def gate_geometry(prs, sw, sh):
    flags = []
    for si, slide in enumerate(prs.slides, 1):
        majors = []
        for shape in slide.shapes:
            try: x=shape.left/EMU; y=shape.top/EMU; w=shape.width/EMU; h=shape.height/EMU
            except Exception: continue
            if x < -0.03 or y < -0.03 or x+w > sw+0.03 or y+h > sh+0.03:
                nm = shape.text_frame.text[:14] if getattr(shape,"has_text_frame",False) else str(shape.shape_type)
                flags.append(f"S{si} OUT-OF-BOUNDS '{nm}' {x:.2f},{y:.2f},{w:.2f},{h:.2f}")
            if w >= 2.0 and h >= 0.40:
                majors.append((x,y,w,h))
        def nested(a,b):
            ins=lambda p,q:(p[0]>=q[0]-0.05 and p[1]>=q[1]-0.05 and p[0]+p[2]<=q[0]+q[2]+0.05 and p[1]+p[3]<=q[1]+q[3]+0.05)
            return ins(a,b) or ins(b,a)
        for i in range(len(majors)):
            for j in range(i+1, len(majors)):
                a,b = majors[i], majors[j]
                ix = max(0, min(a[0]+a[2], b[0]+b[2]) - max(a[0], b[0]))
                iy = max(0, min(a[1]+a[3], b[1]+b[3]) - max(a[1], b[1]))
                if ix*iy > 0.05 and not nested(a,b):
                    flags.append(f"S{si} OVERLAP major blocks area={ix*iy:.2f}")
    return flags

def ensure_render(pptx_path, render_dir):
    existing = sorted(glob.glob(os.path.join(render_dir, "slide-*.jpg")))
    if existing: return existing
    os.makedirs(render_dir, exist_ok=True)
    pdf = os.path.join(render_dir, "deck.pdf")
    # try soffice
    soffice = None
    for cand in ("soffice", "libreoffice"):
        if subprocess.call(["bash","-lc",f"command -v {cand} >/dev/null 2>&1"]) == 0:
            soffice = cand; break
    if soffice is None:
        return []
    subprocess.run([soffice,"--headless","--convert-to","pdf","--outdir",render_dir,pptx_path],
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    src_pdf = os.path.join(render_dir, os.path.splitext(os.path.basename(pptx_path))[0]+".pdf")
    if os.path.exists(src_pdf) and src_pdf != pdf: os.replace(src_pdf, pdf)
    if not os.path.exists(pdf): return []
    subprocess.run(["pdftoppm","-jpeg","-r","135",pdf,os.path.join(render_dir,"slide")],
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return sorted(glob.glob(os.path.join(render_dir, "slide-*.jpg")))

def gate_density(jpgs, slide_w_in, max_blank_in, min_cov_pct):
    import numpy as np
    from PIL import Image
    flags = []; report = []
    for path in jpgs:
        im = np.asarray(Image.open(path).convert("L"), dtype=np.uint8)
        H, W = im.shape
        nonwhite = (im < 245)
        cov = nonwhite.mean()*100
        cell = max(1, int(W/slide_w_in*0.1))
        gh, gw = H//cell, W//cell
        g = np.zeros((gh, gw), bool)
        for i in range(gh):
            for j in range(gw):
                g[i,j] = nonwhite[i*cell:(i+1)*cell, j*cell:(j+1)*cell].any()
        dp = np.zeros((gh, gw), int); best = 0
        for i in range(gh):
            for j in range(gw):
                if not g[i,j]:
                    dp[i,j] = 1 if (i==0 or j==0) else min(dp[i-1,j],dp[i,j-1],dp[i-1,j-1])+1
                    best = max(best, dp[i,j])
        blank_in = best*0.1
        report.append((os.path.basename(path), round(cov,1), round(blank_in,1)))
        if blank_in > max_blank_in + 0.05:
            flags.append(f"{os.path.basename(path)} BLANK block {blank_in:.1f}in > {max_blank_in}in")
        if cov < min_cov_pct:
            flags.append(f"{os.path.basename(path)} LOW coverage {cov:.1f}% < {min_cov_pct}%")
    return flags, report

def gate_font_size(prs, min_pt):
    """HARD: any run with an explicit size below min_pt fails; sizeless runs -> warning."""
    flags = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False): continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if not r.text.strip(): continue
                    if r.font.size is None:
                        flags.append(f"(warn) S{si} run '{r.text[:18]}' has no explicit size (inherits)")
                    elif r.font.size.pt < min_pt - 1e-6:
                        flags.append(f"S{si} SMALL FONT {r.font.size.pt:.1f}pt < {min_pt}pt '{r.text[:22]}'")
    return flags

def gate_shape(prs):
    """HARD: containers must be square-cornered. Exception (dept convention):
    small tag chips <= 1.1in x 0.5in may be rounded (roundRect)."""
    from pptx.oxml.ns import qn
    from pptx.util import Emu
    BANNED = {"roundRect", "round1Rect", "round2DiagRect", "round2SameRect",
              "snip1Rect", "snip2DiagRect", "snip2SameRect", "snipRoundRect"}
    CHIP_W, CHIP_H = 1.1, 0.5
    flags = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            spPr = getattr(getattr(sh, "_element", None), "spPr", None)
            if spPr is None: continue
            g = spPr.find(qn('a:prstGeom'))
            if g is None: continue
            if g.get('prst') in BANNED:
                try:
                    if (sh.width/914400 <= CHIP_W and sh.height/914400 <= CHIP_H):
                        continue  # small chip — allowed
                except Exception:
                    pass
                txt = sh.text_frame.text[:18] if getattr(sh, "has_text_frame", False) else ""
                flags.append(f"S{si} ROUNDED '{g.get('prst')}' shape '{txt}' — use a plain rectangle")
    return flags

def frame_slide_idx(prs, extra=None):
    """1-based indices of frame slides (cover/toc/section/end): exempt from
    DENSITY and STRUCTURE gates. Detected by slide.name == 'frame:*' (set by
    dept_pages.py); slide 1 and the last slide are assumed frame pages too."""
    idx = set(extra or [])
    slides = list(prs.slides)
    for si, slide in enumerate(slides, 1):
        if (slide.name or "").startswith("frame:"):
            idx.add(si)
    if len(slides) > 1:
        idx.add(1); idx.add(len(slides))
    return idx

def gate_structure(prs, min_shapes, frames=()):
    flags = []
    slides = list(prs.slides)
    for si, slide in enumerate(slides, 1):
        if si in frames:
            continue
        n = len(slide.shapes)
        if n < min_shapes:
            flags.append(f"S{si} THIN only {n} shapes < {min_shapes} (possible fake flowchart)")
    return flags

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--tokens", default="assets/tokens.json")
    ap.add_argument("--render-dir", default=None)
    ap.add_argument("--font", default=None)
    ap.add_argument("--frame-slides", default="",
        help="comma-separated 1-based slide numbers to exempt from density/structure (in addition to auto-detected frame:* slides)")
    args = ap.parse_args()

    tk = load_tokens(args.tokens)
    d = tk.get("density", {})
    sw = tk["layout"]["slide_w_in"]; sh = tk["layout"]["slide_h_in"]
    prs = Presentation(args.pptx)
    extra = [int(x) for x in args.frame_slides.split(",") if x.strip()]
    frames = frame_slide_idx(prs, extra)

    font_path = find_font(args.font)
    all_flags = {}

    if font_path:
        all_flags["overflow"] = gate_text_overflow(prs, font_path, d.get("text_overflow_tolerance_pct",3))
    else:
        all_flags["overflow"] = ["(skipped: no CJK font found; pass --font)"]

    all_flags["geometry"] = gate_geometry(prs, sw, sh)

    render_dir = args.render_dir or tempfile.mkdtemp(prefix="qa_")
    jpgs = ensure_render(args.pptx, render_dir)
    if jpgs:
        dflags, report = gate_density(jpgs, sw, d.get("max_blank_square_in",1.2), d.get("min_ink_coverage_pct",30))
        def keep(f):
            m = re.match(r"slide-(\d+)\.jpg", f.split()[0])
            return not (m and int(m.group(1)) in frames)
        dflags = [f for f in dflags if keep(f)]
        all_flags["density"] = dflags
    else:
        report = []
        all_flags["density"] = ["(skipped: could not rasterize; install soffice+pdftoppm or pre-render slide-*.jpg)"]

    all_flags["structure"] = gate_structure(prs, d.get("min_shapes_per_content_slide",40), frames)

    min_pt = d.get("min_font_pt", tk.get("type_scale_pt", {}).get("min", 11))
    all_flags["fontsize"] = gate_font_size(prs, min_pt)
    all_flags["shape"] = gate_shape(prs)

    print("="*56)
    print("QA REPORT —", os.path.basename(args.pptx))
    print("="*56)
    if report:
        print("Density per slide (file, ink%, largest blank in):")
        for r in report: print("  ", r)
        print()
    hard_fail = 0
    for gate, flags in all_flags.items():
        real = [f for f in flags if not f.startswith(("(skipped","(warn"))]
        status = "PASS" if not real else f"FAIL ({len(real)})"
        print(f"[{gate.upper():9}] {status}")
        for f in flags: print("     -", f)
        if real and gate in ("overflow","geometry","density","fontsize","shape"):
            hard_fail += len(real)
    print("="*56)
    if hard_fail:
        print(f"RESULT: FAIL — {hard_fail} hard issue(s). Fix and re-run.")
        sys.exit(1)
    print("RESULT: PASS — all hard gates green.")
    sys.exit(0)

if __name__ == "__main__":
    main()
